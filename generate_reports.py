"""
Generate Laporan Bulanan April 2026 - DOCX & PPTX
"""
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
import os

OUTPUT_DIR = r"d:\MAGANG MAHIRA TOUR\KERJA MAHIRA TOUR\LAPORAN APRIL"

# ============================================================
# 1. LAPORAN BULANAN APRIL - DOCX
# ============================================================
def create_laporan_bulanan():
    doc = Document()
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('LAPORAN BULANAN KARYAWAN')
    run.bold = True
    run.font.size = Pt(16)
    run.font.name = 'Times New Roman'

    # Info
    info_data = [
        ('NAMA KARYAWAN', 'KHOIRUL GUNAWAN'),
        ('JABATAN', 'IT'),
        ('PERIODE', 'APRIL 2026'),
    ]
    for label, value in info_data:
        p = doc.add_paragraph()
        run = p.add_run(f'{label}\t\t\t: {value}')
        run.font.name = 'Times New Roman'
        run.font.size = Pt(12)

    doc.add_paragraph()

    # Table Kegiatan & Capaian
    table = doc.add_table(rows=1, cols=2)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    # Header
    hdr = table.rows[0].cells
    for i, text in enumerate(['KEGIATAN', 'CAPAIAN']):
        hdr[i].text = text
        for paragraph in hdr[i].paragraphs:
            for run in paragraph.runs:
                run.bold = True
                run.font.name = 'Times New Roman'
                run.font.size = Pt(11)

    # Data rows
    data = [
        ('Finalisasi Aplikasi Sautul Haram\n(Radio Muthawwif)', 
         'Implementasi fitur notifikasi sesi baru berhasil dilakukan.\n'
         'Konten doa, shalawat, dan surah (termasuk Surah Yasin) berhasil ditambahkan dan disempurnakan.\n'
         'Setup offline storage dan halaman riwayat sesi muthawwif berhasil diimplementasikan.\n'
         'Perbaikan bug tampilan beranda, login muthawwif, alur logika sesi, dan UI/UX.\n'
         'Testing di device fisik Android berhasil, APK release tersedia (~100MB).\n'
         'Redesain komponen UI/UX, 7+ improvement, konfigurasi standar keamanan.\n'
         'Optimasi performa aplikasi, font arab, tombol navigasi, dan corak islami.'),

        ('Maintenance & Pengembangan\nWebsite Mahira Tour',
         'Maintenance rutin website berjalan lancar tanpa kendala berarti.\n'
         'Update dokumentasi dan konten website secara berkala.\n'
         'Publish artikel: tour Singapura-Malaysia, testimoni jamaah, lowongan finance.\n'
         'Pembuatan PIC website agar bisa digunakan dan dipelajari oleh staff.\n'
         'Pembayaran dan pengecekan server berhasil dilakukan.'),

        ('Perancangan Proyek\nSistem Laporan Terintegrasi',
         'Riset mendalam proyek terakhir magang, evaluasi opsi proyek.\n'
         'Blueprint dan flowchart sistem berhasil dirancang.\n'
         'Setup awal proyek dan database (Next.js + Supabase) dimulai.\n'
         'Konsep alur kerja harian, mingguan, dan bulanan sudah dirancang.'),

        ('Tugas Operasional',
         'Mengambil 2 spanduk Ust. Nadirman dari mersi dan dibawa ke kantor.\n'
         'Mengambil brosur dari mersi, packing, dan kirim via JNT.'),
    ]

    for kegiatan, capaian in data:
        row = table.add_row().cells
        row[0].text = kegiatan
        row[1].text = capaian
        for cell in row:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(10)

    # Set column widths
    for row in table.rows:
        row.cells[0].width = Cm(5.5)
        row.cells[1].width = Cm(12)

    doc.add_paragraph()
    doc.add_paragraph()

    # Signature
    p = doc.add_paragraph()
    run = p.add_run('Sungai penuh,  1 Mei 2026')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(12)

    p = doc.add_paragraph()
    run = p.add_run('Disetujui oleh:\t\t\t\t\t\t\t\t\tHormat Saya,')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(12)

    for _ in range(3):
        doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run('Khilal, S.Th.I\t\t\tNadirman\t\t\t\t\t\t\tKhoirul Gunawan')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(12)

    p = doc.add_paragraph()
    run = p.add_run('Direktur Utama\t\tKomisaris Utama\t\t\t\t\t\tIT')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(12)

    filepath = os.path.join(OUTPUT_DIR, 'LAPORAN BULAN APRIL.docx')
    doc.save(filepath)
    print(f'[OK] Saved: {filepath}')


# ============================================================
# 2. TARGET KERJA & PERBAIKAN - DOCX
# ============================================================
def create_target_perbaikan():
    doc = Document()
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('TARGET KERJA & PERBAIKAN ')
    run.bold = True
    run.font.size = Pt(14)
    run.font.name = 'Times New Roman'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('Staff IT / FullStack Developer')
    run.font.name = 'Times New Roman'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('Khoirul Gunawan')
    run.font.name = 'Times New Roman'

    doc.add_paragraph()

    # Section 1
    p = doc.add_paragraph()
    run = p.add_run('1. Target Pengembangan Sistem Laporan Terintegrasi (Proyek Terakhir)')
    run.bold = True
    run.font.name = 'Times New Roman'

    targets_1 = [
        'Menyelesaikan setup fondasi proyek (Auth, Role-Based Access, Manajemen User & Divisi)',
        'Mengimplementasikan fitur Laporan Harian (input task pagi, update status sore, submit)',
        'Membangun Dashboard Divisi dan Dashboard Pimpinan untuk monitoring real-time',
        'Mengimplementasikan fitur Pengumuman dan Arsip Dokumen per divisi',
        'Mengintegrasikan Notifikasi Telegram untuk reminder submit laporan',
        'Melakukan testing & UAT bersama staff dan pimpinan',
        'Target penyelesaian: 4 minggu (Mei 2026)',
    ]
    for t in targets_1:
        doc.add_paragraph(t, style='List Bullet')

    doc.add_paragraph()

    # Section 2
    p = doc.add_paragraph()
    run = p.add_run('2. Target Maintenance Proyek yang Sudah Selesai')
    run.bold = True
    run.font.name = 'Times New Roman'

    targets_2 = [
        'Melakukan maintenance rutin Website Mahira Tour (update konten, cek performa)',
        'Monitoring Aplikasi Sautul Haram jika ada feedback atau bug report dari user',
        'Update dokumentasi dan artikel di website secara berkala',
        'Memastikan server dan hosting tetap stabil',
    ]
    for t in targets_2:
        doc.add_paragraph(t, style='List Bullet')

    doc.add_paragraph()

    # Section 3
    p = doc.add_paragraph()
    run = p.add_run('3. Target Dukungan Operasional')
    run.bold = True
    run.font.name = 'Times New Roman'

    targets_3 = [
        'Membantu kebutuhan teknis operasional kantor jika diperlukan',
        'Mendukung kegiatan perusahaan dari sisi IT',
        'Siap membantu pekerjaan lain yang ditugaskan',
    ]
    for t in targets_3:
        doc.add_paragraph(t, style='List Bullet')

    doc.add_paragraph()

    # Perbaikan Kerja
    p = doc.add_paragraph()
    run = p.add_run('PERBAIKAN KERJA')
    run.bold = True
    run.font.size = Pt(14)
    run.font.name = 'Times New Roman'

    perbaikan = [
        'Manajemen waktu lebih terstruktur antara development proyek baru dan maintenance proyek lama',
        'Dokumentasi progress lebih konsisten dan teratur',
        'Testing lebih terencana — siapkan device dan checklist dari awal',
        'Resolusi bug lebih cepat dengan root cause analysis yang lebih baik',
        'Komunikasi progress lebih proaktif kepada atasan',
        'Fokus pada prioritas, kurangi context switching antar task',
    ]
    for pb in perbaikan:
        doc.add_paragraph(pb, style='List Bullet')

    filepath = os.path.join(OUTPUT_DIR, 'TARGET KERJA & PERBAIKAN.docx')
    doc.save(filepath)
    print(f'[OK] Saved: {filepath}')


# ============================================================
# 3. PPT GAMBARAN KASAR - PPTX
# ============================================================
def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(title_text, content_lines, layout_idx=1):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x5E)

        # Content
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.font.name = 'Calibri'
                p.space_after = Pt(6)
        return slide

    # Slide 1: Cover
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Laporan Kinerja IT\nApril 2026'
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = 'Khoirul Gunawan — Staff IT / FullStack Developer\nMahira Tour'

    # Slide 2: Agenda
    add_slide('Agenda Presentasi', [
        '1. Ringkasan Capaian Bulan April',
        '2. Status Proyek Sautul Haram (Selesai)',
        '3. Status Website Mahira Tour (Maintenance)',
        '4. Pengenalan Proyek Terakhir: Sistem Laporan Terintegrasi',
        '5. Target & Rencana Bulan Mei',
        '6. Perbaikan Kerja',
    ])

    # Slide 3: Ringkasan
    add_slide('Ringkasan Capaian Bulan April', [
        'Fokus pada 3 hal: Finalisasi App, Maintenance Web, Persiapan Proyek Baru',
        'Total ~20 hari kerja produktif',
        '15+ bug berhasil diperbaiki di Sautul Haram',
        '3 artikel dipublish di website',
        'Blueprint proyek terakhir selesai dirancang',
        'Setup awal proyek Sistem Laporan Terintegrasi dimulai',
    ])

    # Slide 4: Sautul Haram
    add_slide('Aplikasi Sautul Haram — SELESAI [OK]', [
        'Aplikasi radio muthawwif untuk jamaah umrah sudah selesai dikembangkan',
        '',
        'Fitur utama:',
        '• Audio streaming real-time (muthawwif → jamaah)',
        '• Notifikasi sesi baru',
        '• Panduan doa, shalawat & surah lengkap',
        '• Build APK release tersedia (~100MB)',
        '',
        'Testing berhasil di device Android, mode internet berfungsi baik',
        'Konfigurasi standar keamanan sudah diterapkan',
    ])

    # Slide 5: Website
    add_slide('Website Mahira Tour — MAINTENANCE [OK]', [
        'Website berjalan stabil tanpa kendala major',
        '',
        'Aktivitas maintenance:',
        '• Update dokumentasi berkala',
        '• Publish 3 artikel (tour SG-MY, testimoni, lowongan)',
        '• Pembayaran server berhasil',
        '• PIC website dibuat untuk staff lain',
    ])

    # Slide 6: Kendala
    add_slide('Kendala Bulan April', [
        '1. Error build APK Android — membutuhkan waktu beberapa hari',
        '2. Bug login muthawwif — koneksi database terputus (sudah diatasi)',
        '3. Mode WiFi lokal belum optimal',
        '4. Keterbatasan device testing',
        '5. Beberapa bug UI/UX muncul berulang',
    ])

    # Slide 7: Proyek Terakhir - Pengenalan
    add_slide('Proyek Terakhir: Sistem Laporan Terintegrasi', [
        'Apa itu?',
        'Web app yang mendigitalkan pelaporan kerja harian, arsip dokumen,',
        'dan komunikasi internal Mahira Tour.',
        '',
        'Masalah saat ini:',
        '• Laporan kerja harian masih manual via WhatsApp',
        '• Tidak ada sistem terpusat untuk arsip dokumen divisi',
        '• Pimpinan kesulitan memantau progres kerja semua divisi',
        '• Tidak ada reminder otomatis untuk submit laporan',
    ])

    # Slide 8: Fitur Sistem
    add_slide('Fitur Utama Sistem Laporan', [
        '• Login & Role — Staff dan Pimpinan punya akses berbeda',
        '• Laporan Harian — Input task pagi, update sore, submit',
        '• Dashboard Pimpinan — Lihat semua divisi, siapa submit/belum',
        '• Dashboard Divisi — Setiap divisi punya halaman sendiri',
        '• Pengumuman — Pimpinan broadcast info ke staff',
        '• Arsip Dokumen — Upload & kelola dokumen per divisi',
        '• Notifikasi Telegram — Reminder jam 16:00',
        '• Izin/Ketidakhadiran — Sistem tahu siapa yang memang tidak submit',
    ])

    # Slide 9: Alur Kerja
    add_slide('Alur Kerja Sistem', [
        'PAGI (07-09)     → Staff input rencana kerja hari ini',
        'SIANG                → Staff bekerja normal (tidak perlu buka sistem)',
        'SORE (16:00)     → Telegram reminder "Jangan lupa submit!"',
        'SORE (16-17)     → Staff update status task & submit laporan',
        'MALAM              → Pimpinan review dashboard & beri catatan',
        '',
        'Akses: Web app via HP dan Laptop (mobile-first)',
        'User: ~15-20 karyawan + 2 pimpinan',
    ])

    # Slide 10: Teknologi & Biaya
    add_slide('Teknologi & Biaya', [
        'Teknologi: Next.js + Supabase + Telegram Bot',
        'Domain: laporan.mahiratour.id (subdomain yang sudah ada)',
        '',
        '💰 Biaya Bulanan: Rp 0 (menggunakan free tier)',
        '',
        'Timeline: 4 minggu (Mei 2026)',
        '  Minggu 1: Fondasi (Login, User, Izin)',
        '  Minggu 2: Core (Laporan Harian)',
        '  Minggu 3: Dashboard (Divisi & Pimpinan)',
        '  Minggu 4: Polish (Arsip, Notifikasi, Testing)',
    ])

    # Slide 11: Target Mei
    add_slide('Target Bulan Mei', [
        '1. Menyelesaikan Sistem Laporan Terintegrasi (4 minggu)',
        '2. Maintenance rutin Website Mahira Tour',
        '3. Monitoring Aplikasi Sautul Haram',
        '4. Mendukung kegiatan operasional kantor',
    ])

    # Slide 12: Perbaikan
    add_slide('Perbaikan Kerja', [
        '1. Manajemen waktu lebih terstruktur',
        '2. Dokumentasi progress lebih konsisten',
        '3. Testing lebih terencana (siapkan device dari awal)',
        '4. Resolusi bug lebih cepat',
        '5. Komunikasi progress lebih proaktif',
        '6. Fokus pada prioritas, kurangi context switching',
    ])

    # Slide 13: Penutup
    add_slide('Terima Kasih', [
        '"Dengan selesainya Sautul Haram dan Website Mahira Tour,',
        'kita sekarang fokus ke proyek terakhir yang akan memberikan',
        'dampak langsung pada efisiensi kerja seluruh tim."',
        '',
        'Terbuka untuk pertanyaan dan masukan.',
    ])

    filepath = os.path.join(OUTPUT_DIR, 'Laporan-Kinerja-IT-April-2026.pptx')
    prs.save(filepath)
    print(f'[OK] Saved: {filepath}')


if __name__ == '__main__':
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print('Generating reports...')
    create_laporan_bulanan()
    create_target_perbaikan()
    create_ppt()
    print('\n[OK] All reports generated successfully!')
    print(f'📁 Output directory: {OUTPUT_DIR}')
